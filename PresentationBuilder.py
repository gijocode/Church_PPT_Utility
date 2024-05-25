import datetime
import json
import logging
import os
import platform
from copy import deepcopy
from enum import Enum

import six
from BibleExtractor import BibleExtractor
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

logger = logging.getLogger("ppt_util")
fh = logging.FileHandler("ppt_log.log")
fh.setLevel(logging.INFO)
formatter = logging.Formatter("%(asctime)s - %(name)s - %(levelname)s - %(message)s")
fh.setFormatter(formatter)
logger.addHandler(fh)
logger.setLevel(logging.INFO)

operating_sys = platform.system()


class Service(Enum):
    MALAYALAM = 1
    ENGLISH = 2
    ENGLISH_ARIAL = 3
    BIBLE_PORTION = 4


class PresentationBuilder:
    TEMPLATE_FOLDER = "templates"
    ASSETS_FOLDER = "assets"
    TEMPLATE_FILE_FORMAT = "{service_type}_service_template.pptx"
    KK_FULL_FILE = "kk_full.json"
    LOG_FILE = "ppt_log.log"

    REQUIRED_SLIDES = (
        "first_slide",
        "communion",
        "doxology",
        "offertory",
        "birthday",
        "gospel",
        "epistle",
        "second_lesson",
        "between_lessons",
        "first_lesson",
        "opening_song",
        "template_bible_verse",
        "template_bible_heading",
        "template_song_verse",
    )
    SONGS = (
        "opening_song",
        "between_lessons",
        "birthday",
        "offertory",
        "communion",
        "doxology",
    )
    PORTIONS = ("first_lesson", "second_lesson", "epistle", "gospel")

    def __init__(self, service_type) -> None:
        self.service_type = service_type
        self.presentation = self._load_template()
        self.template_dict = self._create_template_dict()
        self.kk_dict = self._load_kk_dict()
        self.bible_extractor = BibleExtractor()

    def _load_template(self):
        template_file = os.path.join(
            self.TEMPLATE_FOLDER,
            self.TEMPLATE_FILE_FORMAT.format(
                service_type=self.service_type.name.lower()
            ),
        )
        return Presentation(template_file)

    def _create_template_dict(self):
        return {
            shape.name: slide
            for slide in self.presentation.slides
            for shape in slide.shapes
            if shape.name in self.REQUIRED_SLIDES
        }

    def _load_kk_dict(self):
        kk_file = os.path.join(self.ASSETS_FOLDER, self.KK_FULL_FILE)
        with open(kk_file, "r", encoding="utf-8") as kkfile:
            return json.load(kkfile)

    def duplicate_slide(self, template_slide):
        blank_slide_layout = self.presentation.slide_layouts[
            len(self.presentation.slide_layouts) - 1
        ]
        copied_slide = self.presentation.slides.add_slide(blank_slide_layout)

        for shape in template_slide.shapes:
            el = shape.element
            copied_slide.shapes._spTree.insert_element_before(deepcopy(el), "p:extLst")

        for _, value in six.iteritems(template_slide.part.rels):
            if "notesSlide" not in value.reltype:
                copied_slide.part.rels.add_relationship(
                    value.reltype, value._target, value.rId
                )

        return copied_slide

    @staticmethod
    def _get_next_sunday():
        return datetime.date.today() + datetime.timedelta(
            days=(6 - datetime.date.today().weekday() + 7) % 7
        )

    @property
    def xml_slides(self):
        return self.presentation.slides._sldIdLst  # pylint: disable=protected-access

    def _update_first_slide(self, topic):
        first_slide = self.template_dict["first_slide"]
        next_sunday = self._get_next_sunday()
        for shape in first_slide.shapes:
            if shape.name not in ["theme", "todays_date"]:
                continue
            tf = shape.text_frame
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for para in tf.paragraphs:
                para.text = (
                    topic
                    if shape.name == "theme"
                    else next_sunday.strftime("%-d %B, %Y")
                )
                para.font.name = (
                    "Goudy Bookletter 1911" if shape.name == "theme" else "Arial"
                )
                para.font.size = Pt(50)
                para.font.color.rgb = (
                    RGBColor(0, 13, 200)
                    if shape.name == "theme"
                    else RGBColor(235, 114, 8)
                )
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER

    def put_in_ppt(
        self,
        content_list: list[str],
        template_name,
        after_slide,
        line_spacing: float = 1,
    ):
        content_slides = []
        for content in content_list[::-1]:
            verse_slide = self.duplicate_slide(self.template_dict[template_name])
            for shape in verse_slide.shapes:
                if shape.name != template_name:
                    continue
                tf = shape.text_frame
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                tf.clear()
                for para in tf.paragraphs:
                    para.text = content
                    if content.count("\n") < 6:
                        para.line_spacing = line_spacing
                    if template_name == "template_bible_heading":
                        para.font.name = "Noto Serif Malayalam"
                        para.font.size = Pt(54)
                        para.font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        para.font.size = Pt(40)
                        para.font.name = "Goudy Bookletter 1911"
                    para.font.bold = True
                    defrpr = para._element.pPr.defRPr
                    ea = etree.SubElement(defrpr, qn("a:cs"))
                    ea.set("typeface", "Noto Serif Malayalam")
                    para.alignment = PP_ALIGN.CENTER

            content_slides.append(verse_slide)
        index_of_slide = (
            self.presentation.slides.index(self.template_dict[after_slide])
            if after_slide
            else 0
        )
        for i in range(len(content_slides)):
            self.move_slide(-1, index_of_slide + 1 + i)

    def get_kk_lyrics(self, song_no):
        lyrics = self.kk_dict[song_no]["Song"]
        lyrics_slide = lyrics.split("\n\n")
        if lyrics_slide[1].startswith("1"):
            chorus = lyrics_slide[0]
            for i in range(2, 2 * len(lyrics_slide), 2):
                lyrics_slide.insert(i, chorus)
        song_name = self.kk_dict[song_no]["Name"]
        print(song_name)
        title = f"Song No: {self.kk_dict[song_no]['No']}\n{song_name.split('(')[0]}"
        if author := self.kk_dict[song_no].get("Author"):
            title += f"\n\nLyricist: {author}"
        lyrics_slide.insert(0, title)
        return lyrics_slide

    def get_bible_portions(self, portion_type):
        print("\033c", end="")
        print(portion_type.replace("_", " ").upper())
        (
            book,
            chapter,
            starting_verse,
            ending_verse,
        ) = self.bible_extractor.get_input_from_user(portion_type)

        book += (
            39
            if portion_type in {"second_lesson", "gospel"}
            else 43 if portion_type == "epistle" else 0
        )

        if portion_type == "all":
            portion_type = ""

        bible_portion = [
            "{5}\n\n{4} {1}: {2}-{3}\n{0} {1}: {2}-{3}".format(
                self.bible_extractor.bible_books_eng[book],
                chapter + 1,
                starting_verse + 1,
                ending_verse + 1,
                self.bible_extractor.bible_books_mal[book],
                portion_type.upper().replace("_", " "),
            )
        ]
        if not portion_type:
            bible_portion[0] = bible_portion[0].strip("\n")
        portion = [
            "{}\n\n{}".format(
                self.bible_extractor.extract_bible_portion(
                    self.bible_extractor.malbible, book, chapter, i
                ),
                self.bible_extractor.extract_bible_portion(
                    self.bible_extractor.engbible, book, chapter, i
                ),
            )
            for i in range(starting_verse, ending_verse + 1)
        ]

        logger.info(
            "Retrieved bible portion for {} which is {}", portion_type, bible_portion
        )
        self.put_in_ppt(portion, "template_bible_verse", portion_type, 1)
        self.put_in_ppt(bible_portion, "template_bible_heading", portion_type, 1)

    def get_song(self, song_type):
        songs = (
            x.strip()
            for x in input(f"\nEnter songs for {song_type.replace('_',' ')}: ").split(
                ","
            )[::-1]
        )
        for song in songs:
            if not song.isdigit() and not song.startswith("dox"):
                song_lyrics = [song]
            else:
                song_lyrics = self.get_kk_lyrics(song)
            self.put_in_ppt(song_lyrics, "template_song_verse", song_type, 1.5)

    def move_slide(self, old_index, new_index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[old_index])
        self.xml_slides.insert(new_index, slides[old_index])

    def delete_slide(self, index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[index])

    def save_ppt(self, ppt_name=None):
        next_sunday = (
            datetime.date.today()
            + datetime.timedelta(days=(6 - datetime.date.today().weekday() + 7) % 7)
        ).strftime("%-d %B, %Y")
        if operating_sys == "Darwin":
            ppt_name = (
                (
                    f"/Users/gijomathew/Important/misc/Church/PPTs/2024/{next_sunday}.pptx"
                )
                if not ppt_name
                else ppt_name
            )
            self.presentation.save(ppt_name)
            os.system("open '/Users/gijomathew/Important/misc/Church/PPTs/2024/'")
            os.system(f"open '{ppt_name}'")
        else:
            ppt_name = next_sunday if not ppt_name else ppt_name
            self.presentation.save(ppt_name)
        print(f"PPT Successfully saved to {ppt_name}")
        logger.info(f"PPT Successfully saved to {ppt_name}")


if __name__ == "__main__":
    while True:
        choice = int(
            input(
                "What presentation do you want to create? Enter the index:\n1. Malayalam Service\n2. English Service\n3. English Service (Arial Font) \n4. Bible Portion\n"
            )
        )
        if choice in {1, 2, 3, 4}:
            break

    service_type = Service(choice)
    pb = PresentationBuilder(service_type)

    if service_type == Service.BIBLE_PORTION:
        pb.get_bible_portions("all")
        pb.save_ppt("bible_portion.pptx")
        exit(0)

    _ = [pb.get_bible_portions(portion) for portion in pb.PORTIONS]

    print("\033c", end="")
    print(
        'Note:\n1. If you need multiple songs (eg. for Communion) , separate them with commas.\n2. Use "dox1" for Doxology 1, "dox2" for Doxology 2, and so on.\n3. If the song is not from the book "Kristeeya Keerthanangal", simply type the song title. (Lyrics will not be provided).'
    )
    _ = [pb.get_song(song) for song in pb.SONGS]

    try:
        pb._update_first_slide(input("\nEnter the theme for this Sunday: "))
    except Exception as e:
        print(
            "There was some error encountered when updating the first slide. Please make sure you edit it manually!"
        )
        logger.exception(e)

    pb.save_ppt()
