import collections
import collections.abc

from pptx import Presentation

prs = Presentation("templates/PPT English Service.pptx")
prs.slides[1].name = "Opening Hymn"
prs.slides[2].name = "First Lesson"
prs.slides[3].name = "Song between Lessons"
prs.slides[4].name = "Second Lesson"
prs.slides[15].name = "Epistle"
prs.slides[23].name = "Gospel"
prs.slides[44].name = "Birthday,Wedding"
prs.slides[45].name = "Offertory"
prs.slides[118].name = "Communion Songs"
prs.slides[132].name = "Doxology"

prs.save("templates/PPT English Service.pptx")
