import datetime
import json
import logging
import os
import platform
from collections import abc
from copy import deepcopy

import six
from lxml import etree
from pptx import Presentation, exc
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt
from pptx.dml.color import RGBColor

from BibleExtractor import BibleConverter

logger = logging.getLogger("ppt_util")
fh = logging.FileHandler("ppt_log.log")
fh.setLevel(logging.INFO)
formatter = logging.Formatter("%(asctime)s - %(name)s - %(levelname)s - %(message)s")
fh.setFormatter(formatter)
logger.addHandler(fh)
logger.setLevel(logging.INFO)
operating_sys = platform.system()


class PresentationBuilder(object):
    def __init__(self, service_type) -> None:
        logger.info("Created ppt object")
        logger.info(f"Creating PPT on {datetime.date.today()}")
        template = f"templates/{service_type}_template.pptx"
        self.presentation = Presentation(template)

    def shadow_whole_ppt(self):
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                shape.shadow.inherit=True

    def save_ppt(self, ppt_name=None):
        self.presentation.save(ppt_name)
        os.system(f"open '{ppt_name}'")


if __name__ == "__main__":

    service_type = "malayalam_service"
    pb_obj = PresentationBuilder(service_type)
    pb_obj.shadow_whole_ppt()
    pb_obj.save_ppt("test.pptx")
