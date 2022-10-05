import collections
import collections.abc
import os
from datetime import datetime

from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

from bible_extractor import BibleExtractor

bible_ext = BibleExtractor()
prs = Presentation("templates/PPT English Service.pptx")
slide_contents = {
    "Opening Hymn": "",
    "First Lesson": "",
    "Second Lesson": "",
    "Song between Lessons": "",
    "Epistle": "",
    "Gospel": "",
    "Birthday,Wedding": "",
    "Offertory": "",
    "Communion Songs": "",
    "Doxology": "",
}

slide_map = {slide.name: slide for slide in prs.slides}

# Title slide modifications
ppt_date = input("Enter date in dd-mm-yy format")
theme = input("Enter the theme for this Sunday")
title_slide = slide_map["Title"]
title_slide.shapes[2].text = datetime.strptime(ppt_date, "%d-%m-%y").strftime(
    "%d %B, %Y"
)
title_slide.shapes[2].text_frame.fit_text(font_family="Arial", max_size=100)
title_slide.shapes[2].text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # type: ignore

title_slide.shapes[4].text = f"Theme: {theme}"
title_slide.shapes[4].text_frame.fit_text(font_family="Arial", max_size=100)
title_slide.shapes[4].text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # type: ignore

# Getting First lesson notation and verses
notation, verses = bible_ext.get_first_lesson()
slide_contents["First Lesson"] = notation
with open("output_bible_portions/first_lesson.txt", "w") as first_lesson_file:
    first_lesson_file.write(verses)


# Getting Second lesson notation and verses
notation, verses = bible_ext.get_second_lesson()
slide_contents["Second Lesson"] = notation
with open("output_bible_portions/second_lesson.txt", "w") as second_lesson_file:
    second_lesson_file.write(verses)

# Getting Epistle notation and verses
notation, verses = bible_ext.get_epistle()
slide_contents["Epistle"] = notation
with open("output_bible_portions/epistle.txt", "w") as epistle_file:
    epistle_file.write(verses)

# Getting Gospel notation and verses
notation, verses = bible_ext.get_gospel()
slide_contents["Gospel"] = notation
with open("output_bible_portions/gospel.txt", "w") as gospel_file:
    gospel_file.write(verses)

# Here we will update the songs for various sections

slide_contents["Opening Hymn"] = input("\nEnter opening hymn\n")
slide_contents["Song between Lessons"] = input("\nEnter song between lessons")
slide_contents["Birthday,Wedding"] = input("\nEnter birthday hymn\n")
slide_contents["Offertory"] = input("\nEnter offertory hymn\n")
no_of_communion_songs = int(input("\nEnter no of communion songs: "))
comm_songs = []
for i in range(no_of_communion_songs):
    comm_songs.append(input(f"Enter song no {i+1}\n"))
slide_contents["Communion Songs"] = "\n".join(comm_songs)
slide_contents["Doxology"] = input("\nEnter Doxology Hymn")

# Fill the slides with collected values
for slide, slide_text in slide_contents.items():
    text_element = slide_map[slide].shapes[1]
    text_element.text = slide_text
    text_element.text_frame.fit_text(font_family="Arial", max_size=100)
    for paragraph in text_element.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # type: ignore


os.system("clear")  # If using windows, use 'cls'


# Saving final PPT with updated portions and songs
prs.save("output_ppt/test.pptx")
print("Saved the PPT!")
